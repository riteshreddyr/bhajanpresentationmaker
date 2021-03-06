__author__ = 'RiteshReddy'
import os
import random

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from ttfquery import describe, glyphquery

from flaskappbase import app


class SaiPresentation():
    def __init__(self, title, subtitle, title_silde_image=os.path.join(app.config['DATA_DIRECTORY'], "title_slide.jpg")):
        self.prs = Presentation()
        self.add_title_slide(title, subtitle, title_silde_image)

    @staticmethod
    def set_run_font(run, size, color=RGBColor(0x00, 0x00, 0x00), bold=False, italic=None):
        font = run.font
        font.size = size
        font.color.rgb = color
        font.bold = bold
        font.italic = italic

    @staticmethod
    def add_new_run_with_text(para, text=""):
        run = para.add_run()
        run.text = text
        return run

    @staticmethod
    def add_paragraph_with_alignment(textBox, alignment=PP_ALIGN.CENTER):
        para = textBox.text_frame.add_paragraph()
        para.alignment = alignment
        return para

    @staticmethod
    def add_textBox(slide, left, top, width, height, auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, word_wrap = True):
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box.text_frame.auto_size = auto_size
        text_box.text_frame.word_wrap = word_wrap
        return text_box

    @staticmethod
    def add_title_box_runs(slide, color=RGBColor(0xff, 0xff, 0xff)):
        text_box = SaiPresentation.add_textBox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(1.05))
        para = SaiPresentation.add_paragraph_with_alignment(text_box, PP_ALIGN.LEFT)
        title_run = SaiPresentation.add_new_run_with_text(para, "")
        misc_run = SaiPresentation.add_new_run_with_text(para, "")
        SaiPresentation.set_run_font(title_run, Pt(36), color, False, None)
        SaiPresentation.set_run_font(misc_run, Pt(24), color, False, None)
        return title_run, misc_run

    @staticmethod
    def add_run_to_slide_with_font(slide, left, top, width, height, text_size, alignment=PP_ALIGN.CENTER,
                                   color=RGBColor(0xff, 0xff, 0xff), bold=False, italic=None):
        text_box = SaiPresentation.add_textBox(slide, left, top, width, height)
        para = SaiPresentation.add_paragraph_with_alignment(text_box, alignment)
        run = SaiPresentation.add_new_run_with_text(para, "")
        SaiPresentation.set_run_font(run, text_size, color, bold, italic)
        return run

    def add_full_image_slide(self, img_path):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        left = top = Inches(0)
        height = Inches(7.5)
        width = Inches(10)
        pic = slide.shapes.add_picture(img_path, left, top, height=height, width=width)

    def add_title_slide(self, title, subtitle, image_path="title_slide.jpg"):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        left = top = Inches(0)
        height = Inches(7.5)
        width = Inches(10)
        pic = slide.shapes.add_picture(image_path, left, top, height=height, width=width)
        text_box = SaiPresentation.add_textBox(slide, Inches(2.5), Inches(2.2), Inches(5), Inches(4))
        para = SaiPresentation.add_paragraph_with_alignment(text_box)
        run = SaiPresentation.add_new_run_with_text(para, title)
        SaiPresentation.set_run_font(run, Pt(54), RGBColor(0xFF, 0xFF, 0xFF), True)
        run = SaiPresentation.add_new_run_with_text(para, "\n") # new line.
        run = SaiPresentation.add_new_run_with_text(para, subtitle)
        SaiPresentation.set_run_font(run, Pt(32), RGBColor(0xFF, 0xFF, 0x0F), True)


    def add_bhajan_slide(self, bhajan_name, bhajan_txt, bhajan_gender="", key="", next_bhajan_name="", next_bhajan_gender = "", next_key="", backgroundImage=None, noBackground=None, hexTextColor=None):
        """
        Adds a bhajan to the powerpoint - a bhajan can take multiple slides depending on user
        handled pagebreaks and natural overflows.
        One slide can hold 9 rows and 49 characters per row.
        :param bhajan_name:
        :param bhajan_txt:
        :param key:
        :param next_bhajan_name:
        :param next_key:
        :param backgroundImage:
        :param hexTextColor:
        :return:
        """
        MAX_ROW_COUNT = 9
        SIZE_OF_A = 1185
        MAX_ROW_LENGTH_CALIBRI = 33 * SIZE_OF_A # 33 'A's of size 32 pt
        MAX_TITLE_LENGTH = 24 * SIZE_OF_A # 24 As of size 36 pt
        calibri = describe.openFont(os.path.join(app.config['DATA_DIRECTORY'], "calibri.ttf"))

        def get_character_width(character):
            return glyphquery.width(calibri, glyphquery.glyphName(calibri, character))

        def bhajan_slide_template(background_path=None, color=None):
            """
            Adds a slide to the powerpoint and returns a handle to all the various text boxes
            :param background_path: If not None, include this image as the background image.
            :return: slide, title_box, bhajan_box, key_box, next_bhajan_name_box, next_key_box
            """
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            if not background_path is None:
                slide.shapes.add_picture(background_path, Inches(0), Inches(0), height=Inches(7.5), width=Inches(10))
            title_rn, misc_rn = self.add_title_box_runs(slide, color=color)
            bhajan_rn = self.add_run_to_slide_with_font(slide, Inches(0.5), Inches(1.25), Inches(9), Inches(5.5),
                                                        Pt(32), PP_ALIGN.LEFT, color=color)
            next_label_rn = self.add_run_to_slide_with_font(slide, Inches(0.5), Inches(6.5), Inches(1), Inches(0.75), Pt(16),
                                                     PP_ALIGN.LEFT, color=color)
            next_bhajan_name_rn = self.add_run_to_slide_with_font(slide, Inches(2.5), Inches(6.5), Inches(4),
                                                                  Inches(0.75), Pt(16), color=color)
            next_key_rn = self.add_run_to_slide_with_font(slide, Inches(8), Inches(6.5), Inches(1), Inches(0.75),
                                                          Pt(16), PP_ALIGN.RIGHT, color=color)
            return slide, title_rn, misc_rn, bhajan_rn, next_label_rn, next_bhajan_name_rn, next_key_rn

        def handle_user_defined_page_breaks(text):
            def split_lines_into_slide_using_max_row_count_per_slide(text):
                """
                Takes a bunch of lines and groups them into MAX_ROW_COUNT per slide
                :param text: str - Bhajan text with line breaks to indicate lines.
                :return: list of strings - each element represents the text for a slide.
                """
                def get_true_lines(text):
                    line_list = text.split("\r\n")
                    for ind, line in enumerate(line_list):
                        if len(line) > MAX_ROW_LENGTH_CALIBRI:
                            """ Break into spaces and reconstruct at 35 intervals : handles multiple lines in one go """
                            new_line = "" #string concat...so bad :(
                            cur_line_length = 0
                            for char in line:
                                width = get_character_width(char)
                                if cur_line_length == 0:
                                    new_line += char
                                    cur_line_length += width
                                elif cur_line_length + width + 1 > MAX_ROW_LENGTH_CALIBRI:
                                    new_line += "\n" + char
                                    cur_line_length = width
                                else:
                                    new_line += " " + char
                                    cur_line_length += width + 1
                            line_list[ind] = new_line
                    return "\n".join(line_list)
                slide_text = []
                lines = get_true_lines(text).split("\n")
                this_set = []
                for i in range(0, len(lines), MAX_ROW_COUNT):
                    for j in range(i, min(len(lines), i + MAX_ROW_COUNT)):
                        this_set.append(lines[j])
                    slide_text.append('\n'.join(this_set))
                    this_set = []
                return slide_text

            # A user can define manual page breaks with the following marker.
            # Split the text by the marker and ensure each split is in a different slide.
            PAGE_BREAK_MARKER = '\r\n[pagebreak]\r\n'
            slide_text = []
            for line in text.split(PAGE_BREAK_MARKER):
                slide_text.extend(split_lines_into_slide_using_max_row_count_per_slide(line))
            return slide_text

        def truncate_title(bhajan_name):
            if not bhajan_name:
                return bhajan_name
            prefix_sum = [0] * len(bhajan_name)
            prefix_sum[0] = get_character_width(bhajan_name[0])
            break_at = -1
            for i in range(1, len(bhajan_name)):
                prefix_sum[i] = prefix_sum[i-1] + get_character_width(bhajan_name[i])
                if prefix_sum[i] >= MAX_TITLE_LENGTH:
                    break_at = i
                    break
            if break_at == -1:
                return bhajan_name # full name can fit.
            for i in range(break_at, 0, -1):
                if bhajan_name[i] == " ":
                    return bhajan_name[:i] + "..."
            return bhajan_name # cannot break, return full name.


        def add_a_bhajan_slide(bhajan_name, text, final=True, bhajan_gender="", key="", next_bhajan_name="",
                               next_bhajan_gender="", next_key="", background_path=None, color=None):
            """
            Adds a single bhajan slide - this is one in a series that a single bhajan can take up
            :param bhajan_name:
            :param text:
            :param final: - whether this is the final slide for this bhajan or not.
            :param key:
            :param next_bhajan_name:
            :param next_key:
            :param background_path: - image path for the background picture
            """
            slide, title_rn, misc_rn, bhajan_rn, next_label_rn, nxt_bhajan_rn, nxt_key_rn = bhajan_slide_template(background_path, color)
            title_rn.text = truncate_title(bhajan_name.strip())
            if bhajan_gender:
                misc_rn.text += " (" + bhajan_gender + ")"
            if key:
                misc_rn.text += " - " + key.strip()
            bhajan_rn.text = text.strip()
            if not final:
                nxt_bhajan_rn.text = "Continued"
            else:
                if next_bhajan_name:
                    next_label_rn.text = "Next"
                nxt_bhajan_rn.text = next_bhajan_name
                if next_bhajan_gender:
                    nxt_bhajan_rn.text += " (" + next_bhajan_gender + ")"
                nxt_key_rn.text = next_key

        # background images
        background_path = None
        if not noBackground:
            if backgroundImage:
                background_path = backgroundImage
            else:
                # background images - pick a random one.
                files = os.listdir(os.path.join(app.config['DATA_DIRECTORY'], 'backgrounds'))
                if len(files) != 0:
                    index = random.randint(0, len(files) - 1)
                    background_path = os.path.join(os.path.join(app.config['DATA_DIRECTORY'], 'backgrounds'), files[index])
        rgbColor = RGBColor(0xff, 0xff, 0xff)
        if hexTextColor and len(hexTextColor) == 6:
            r = int(hexTextColor[0:2], 16)
            g = int(hexTextColor[2:4], 16)
            b = int(hexTextColor[4:6], 16)
            rgbColor = RGBColor(r,g,b)

        text_split_per_slide = handle_user_defined_page_breaks(bhajan_txt)
        for pos, text in enumerate(text_split_per_slide, start=1):
            if pos == len(text_split_per_slide):
                final = True
            else:
                final = False
            if len(bhajan_txt) == 0:  # Empty Bhajan
                add_a_bhajan_slide(bhajan_name, text, final, bhajan_gender, key, next_bhajan_name, next_bhajan_gender, next_key, os.path.join(app.config['DATA_DIRECTORY'], 'filler.jpg'), rgbColor)
            else:
                add_a_bhajan_slide(bhajan_name, text, final, bhajan_gender, key, next_bhajan_name, next_bhajan_gender, next_key, background_path, rgbColor)

    def save_presentation(self, filename):
        self.prs.save(filename)

    def create_test_presentation(self):
        self.add_bhajan_slide("Ganesha Sharanam Ganesha Sharanam Ganesha Sharanam Ganesha Sharanam",
                              "Ganesha Ganesha Ganesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha Ganesha\nGanesha Ganesha Ganesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha Ganesha\nGanesha Ganesha Ganesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha Ganesha\nGanesha Ganesha Ganesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha Ganesha\nGanesha Ganesha Ganesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha Ganesha\nGanesha Ganesha Ganesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha GaneshaGanesha Ganesha Ganesha Ganesha Ganesha\n",
                              "Am", "Guru Baba", "B#")
        self.prs.save("test.pptx")


if __name__ == "__main__":
    SaiPresentation().create_test_presentation()